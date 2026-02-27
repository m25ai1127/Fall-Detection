"""Generate PowerPoint presentation for professor meeting (Weeks 1-3)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK     = RGBColor(0x1B, 0x1B, 0x2F)
BG_CARD     = RGBColor(0x27, 0x27, 0x3F)
ACCENT      = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2     = RGBColor(0x90, 0xE0, 0xEF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
GREEN       = RGBColor(0x06, 0xD6, 0xA0)
ORANGE      = RGBColor(0xFF, 0xD1, 0x66)
RED_SOFT    = RGBColor(0xEF, 0x47, 0x6F)


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, size=16, color=LIGHT_GRAY, bold=False,
             bullet=False, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = space_before
    if bullet:
        p.level = 1
    return p


def add_card(slide, left, top, width, height, color=BG_CARD, border=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)

add_accent_line(slide, 1, 2.5, 11.3)

add_text(slide, 1, 2.7, 11.3, 1.2,
         "AI-Based Fall and Abnormal Motion Detection",
         size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, 1, 3.7, 11.3, 0.7,
         "Using Depth-Assisted Single-Camera Vision",
         size=24, color=ACCENT2, align=PP_ALIGN.CENTER)

add_text(slide, 1, 5.0, 11.3, 0.5,
         "Progress Report — Weeks 1 to 3",
         size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 1, 5.6, 11.3, 0.5,
         "February 2026",
         size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 2: Problem Statement
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 6, 0.6, "Problem Statement", size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.8, 0.95, 4)

# Left column
add_card(slide, 0.8, 1.3, 5.5, 5.5)
tf = add_text(slide, 1.1, 1.5, 5, 0.5, "The Challenge", size=20, color=ACCENT2, bold=True)
add_para(tf, "")
add_para(tf, "• Falls are a leading cause of injury among the elderly", size=16, color=WHITE)
add_para(tf, "• Timely detection can reduce serious health complications", size=16, color=WHITE)
add_para(tf, "• Traditional approaches require expensive RGB-D cameras", size=16, color=WHITE)
add_para(tf, "  (Microsoft Kinect, Intel RealSense)", size=14, color=LIGHT_GRAY)
add_para(tf, "")
add_para(tf, "Our Approach", size=20, color=GREEN, bold=True)
add_para(tf, "")
add_para(tf, "• Use a single regular camera (no depth sensor)", size=16, color=WHITE)
add_para(tf, "• Estimate depth using AI (MiDaS monocular depth)", size=16, color=WHITE)
add_para(tf, "• Combine depth + body pose for 3D posture analysis", size=16, color=WHITE)
add_para(tf, "• Classify fall vs normal using temporal AI (LSTM)", size=16, color=WHITE)

# Right column
add_card(slide, 6.8, 1.3, 5.7, 5.5)
tf = add_text(slide, 7.1, 1.5, 5.2, 0.5, "Key Innovation", size=20, color=ACCENT2, bold=True)
add_para(tf, "")
add_para(tf, "Depth-assisted single-camera approach eliminates", size=16, color=WHITE)
add_para(tf, "the need for expensive depth-sensing hardware.", size=16, color=WHITE)
add_para(tf, "")
add_para(tf, "The system works with:", size=16, color=WHITE)
add_para(tf, "  • Any standard webcam or CCTV camera", size=16, color=GREEN)
add_para(tf, "  • Pre-recorded video files", size=16, color=GREEN)
add_para(tf, "  • No special hardware required", size=16, color=GREEN)
add_para(tf, "")
add_para(tf, "Objective", size=20, color=ORANGE, bold=True)
add_para(tf, "")
add_para(tf, "Build a real-time fall detection system that is", size=16, color=WHITE)
add_para(tf, "affordable, accurate, and deployable anywhere.", size=16, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: System Architecture
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 6, 0.6, "System Architecture", size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.8, 0.95, 4)

# Pipeline boxes
boxes = [
    ("RGB Video\nFrame",         1.0, 2.5, 1.8, 1.2, ACCENT),
    ("MiDaS\nDepth Estimation",  3.3, 1.5, 2.0, 1.2, RGBColor(0x48, 0xCA, 0xE4)),
    ("Pose\nEstimation",         3.3, 3.5, 2.0, 1.2, RGBColor(0x48, 0xCA, 0xE4)),
    ("Feature\nExtraction\n(177-dim)", 6.0, 2.5, 2.0, 1.4, GREEN),
    ("LSTM\nClassifier",         8.7, 2.5, 2.0, 1.2, ORANGE),
    ("Fall / Normal\nDecision",  11.2, 2.5, 1.8, 1.2, RED_SOFT),
]

for text, l, t, w, h, clr in boxes:
    shape = add_card(slide, l, t, w, h, color=clr)
    shape.text_frame.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            shape.text_frame.paragraphs[0].text = line
            shape.text_frame.paragraphs[0].font.size = Pt(13)
            shape.text_frame.paragraphs[0].font.color.rgb = BG_DARK
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        else:
            p = shape.text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
            p.font.color.rgb = BG_DARK
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(8)

# Arrows
for x in [2.85, 5.4, 8.1, 10.75]:
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x), Inches(2.9), Inches(0.45), Inches(0.35)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LIGHT_GRAY
    arrow.line.fill.background()

# Branch arrows from RGB to MiDaS and Pose
for y_target in [1.9, 3.9]:
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(2.85), Inches(y_target), Inches(0.45), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LIGHT_GRAY
    arrow.line.fill.background()

# Description below
add_card(slide, 0.8, 4.8, 11.7, 2.2)
tf = add_text(slide, 1.1, 4.9, 11.2, 0.4, "Pipeline Description", size=18, color=ACCENT2, bold=True)
add_para(tf, "1. MiDaS generates a depth map from a regular RGB frame — estimates distance of each pixel from camera", size=14, color=WHITE)
add_para(tf, "2. Pose Estimation detects 33 skeletal joints (head, shoulders, hips, knees, ankles, etc.)", size=14, color=WHITE)
add_para(tf, "3. Feature Extraction combines depth + pose into 177-dimensional vector per frame", size=14, color=WHITE)
add_para(tf, "4. Bidirectional LSTM processes 30 consecutive frames and classifies: Fall or Normal", size=14, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 4: Dataset
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 6, 0.6, "Dataset", size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.8, 0.95, 4)

# URFD card
add_card(slide, 0.8, 1.3, 5.8, 5.5, border=ACCENT)
tf = add_text(slide, 1.1, 1.5, 5.3, 0.5,
              "UR Fall Detection Dataset (URFD)", size=22, color=WHITE, bold=True)
add_para(tf, "Primary Dataset", size=14, color=GREEN, bold=True)
add_para(tf, "")
add_para(tf, "• 70 video sequences total", size=16, color=WHITE)
add_para(tf, "  - 30 fall sequences", size=15, color=ACCENT2)
add_para(tf, "  - 40 ADL (normal activity) sequences", size=15, color=ACCENT2)
add_para(tf, "• Recorded with Microsoft Kinect sensor", size=16, color=WHITE)
add_para(tf, "• Provides both RGB and depth streams", size=16, color=WHITE)
add_para(tf, "• Resolution: 640 x 480", size=16, color=WHITE)
add_para(tf, "")
add_para(tf, "Reference:", size=14, color=LIGHT_GRAY)
add_para(tf, "Kwolek & Kepski, Computers in Biology", size=14, color=LIGHT_GRAY)
add_para(tf, "and Medicine, 2014", size=14, color=LIGHT_GRAY)

# NTU card
add_card(slide, 7.0, 1.3, 5.5, 5.5, border=ORANGE)
tf = add_text(slide, 7.3, 1.5, 5.0, 0.5,
              "NTU RGB+D Dataset", size=22, color=WHITE, bold=True)
add_para(tf, "Supplementary Dataset", size=14, color=ORANGE, bold=True)
add_para(tf, "")
add_para(tf, "• 56,880 action sequences", size=16, color=WHITE)
add_para(tf, "• 60 action classes", size=16, color=WHITE)
add_para(tf, "• Used selectively:", size=16, color=WHITE)
add_para(tf, "  - A043: Falling down", size=15, color=ACCENT2)
add_para(tf, "  - Normal activities for supplementary", size=15, color=ACCENT2)
add_para(tf, "    training data", size=15, color=ACCENT2)
add_para(tf, "")
add_para(tf, "Purpose: Increase diversity of normal", size=16, color=WHITE)
add_para(tf, "motion patterns and improve model", size=16, color=WHITE)
add_para(tf, "generalization.", size=16, color=WHITE)
add_para(tf, "")
add_para(tf, "Reference:", size=14, color=LIGHT_GRAY)
add_para(tf, "Shahroudy et al., CVPR 2016", size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 5: MiDaS Depth Estimation
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 8, 0.6, "Monocular Depth Estimation — MiDaS",
         size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.8, 0.95, 6)

add_card(slide, 0.8, 1.3, 5.8, 5.5)
tf = add_text(slide, 1.1, 1.5, 5.3, 0.5,
              "How It Works", size=22, color=ACCENT2, bold=True)
add_para(tf, "")
add_para(tf, "MiDaS (by Intel ISL) is a pre-trained neural", size=16, color=WHITE)
add_para(tf, "network for monocular depth estimation.", size=16, color=WHITE)
add_para(tf, "")
add_para(tf, "• Takes a single RGB image as input", size=16, color=WHITE)
add_para(tf, "• Outputs a depth map (same resolution)", size=16, color=WHITE)
add_para(tf, "• Darker = closer, Lighter = farther", size=16, color=WHITE)
add_para(tf, "")
add_para(tf, "Model Variant: MiDaS_small", size=16, color=GREEN, bold=True)
add_para(tf, "• Lightweight — suitable for real-time", size=16, color=WHITE)
add_para(tf, "• Based on EfficientNet-Lite3 backbone", size=16, color=WHITE)
add_para(tf, "• Pre-trained on mixed datasets (10+)", size=16, color=WHITE)

add_card(slide, 7.0, 1.3, 5.5, 5.5)
tf = add_text(slide, 7.3, 1.5, 5.0, 0.5,
              "Why Monocular Depth?", size=22, color=ACCENT2, bold=True)
add_para(tf, "")
add_para(tf, "Advantages over hardware depth sensors:", size=16, color=WHITE)
add_para(tf, "")
add_para(tf, "1. No special camera needed", size=16, color=GREEN)
add_para(tf, "   Works with any RGB camera/webcam", size=14, color=LIGHT_GRAY)
add_para(tf, "")
add_para(tf, "2. Lower cost", size=16, color=GREEN)
add_para(tf, "   Kinect: ~$200, RealSense: ~$300", size=14, color=LIGHT_GRAY)
add_para(tf, "   Regular webcam: ~$20", size=14, color=LIGHT_GRAY)
add_para(tf, "")
add_para(tf, "3. Works with existing CCTV footage", size=16, color=GREEN)
add_para(tf, "   Can retrofit existing surveillance systems", size=14, color=LIGHT_GRAY)
add_para(tf, "")
add_para(tf, "Preprocessing (Section 4.2):", size=16, color=ORANGE, bold=True)
add_para(tf, "• Depth normalization to [0, 1]", size=15, color=WHITE)
add_para(tf, "• Noise filtering for missing values", size=15, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 6: Pose Estimation
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 8, 0.6, "Human Pose Estimation",
         size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.8, 0.95, 4)

add_card(slide, 0.8, 1.3, 5.8, 5.5)
tf = add_text(slide, 1.1, 1.5, 5.3, 0.5,
              "Skeletal Landmark Detection", size=22, color=ACCENT2, bold=True)
add_para(tf, "")
add_para(tf, "Extracts 33 key body landmarks per frame", size=16, color=WHITE)
add_para(tf, "")
add_para(tf, "Each landmark provides:", size=16, color=WHITE)
add_para(tf, "  • x, y — position in image (normalized 0-1)", size=15, color=ACCENT2)
add_para(tf, "  • z — relative depth", size=15, color=ACCENT2)
add_para(tf, "  • visibility — detection confidence", size=15, color=ACCENT2)
add_para(tf, "")
add_para(tf, "Key joints for fall detection:", size=16, color=GREEN, bold=True)
add_para(tf, "  • Head (nose) — tracks head height", size=15, color=WHITE)
add_para(tf, "  • Shoulders — body orientation", size=15, color=WHITE)
add_para(tf, "  • Hips — center of gravity", size=15, color=WHITE)
add_para(tf, "  • Knees — leg posture", size=15, color=WHITE)
add_para(tf, "  • Ankles — ground contact", size=15, color=WHITE)

add_card(slide, 7.0, 1.3, 5.5, 5.5)
tf = add_text(slide, 7.3, 1.5, 5.0, 0.5,
              "Why Pose Estimation?", size=22, color=ACCENT2, bold=True)
add_para(tf, "")
add_para(tf, "Falls have characteristic patterns:", size=16, color=WHITE)
add_para(tf, "")
add_para(tf, "Standing", size=16, color=GREEN, bold=True)
add_para(tf, "  Vertical body, joints aligned", size=14, color=LIGHT_GRAY)
add_para(tf, "  Bounding box: tall and narrow", size=14, color=LIGHT_GRAY)
add_para(tf, "")
add_para(tf, "Falling", size=16, color=ORANGE, bold=True)
add_para(tf, "  Rapid downward joint movement", size=14, color=LIGHT_GRAY)
add_para(tf, "  Hip and head drop suddenly", size=14, color=LIGHT_GRAY)
add_para(tf, "")
add_para(tf, "Fallen", size=16, color=RED_SOFT, bold=True)
add_para(tf, "  Horizontal body on ground", size=14, color=LIGHT_GRAY)
add_para(tf, "  Bounding box: wide and short", size=14, color=LIGHT_GRAY)
add_para(tf, "")
add_para(tf, "By tracking joints over time,", size=16, color=WHITE)
add_para(tf, "we capture the full fall motion.", size=16, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 7: Feature Extraction
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 8, 0.6, "Feature Extraction (Section 4.3)",
         size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.8, 0.95, 5)

add_text(slide, 0.8, 1.2, 11.5, 0.5,
         "177-dimensional feature vector per frame combining depth and pose information",
         size=18, color=LIGHT_GRAY)

# Feature table as cards
features = [
    ("Pose Landmarks",           "132", "33 joints x 4 (x, y, z, vis)", ACCENT),
    ("Depth at Joints",          "33",  "Depth value at each joint location", RGBColor(0x48, 0xCA, 0xE4)),
    ("Bounding Box Ratio",       "1",   "Width / Height — posture indicator", GREEN),
    ("Distance to Floor",        "1",   "Height above ground plane", GREEN),
    ("Center of Mass Height",    "1",   "Weighted average joint height", GREEN),
    ("Vertical Velocities",      "9",   "Frame-to-frame downward speed", ORANGE),
]

for i, (name, count, desc, clr) in enumerate(features):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.1
    y = 1.9 + row * 2.6

    add_card(slide, x, y, 3.8, 2.2, border=clr)
    tf = add_text(slide, x + 0.2, y + 0.2, 3.4, 0.4, name, size=17, color=WHITE, bold=True)
    add_para(tf, "")
    add_para(tf, f"Dimensions: {count}", size=22, color=clr, bold=True)
    add_para(tf, desc, size=13, color=LIGHT_GRAY)

# Total
add_card(slide, 4.5, 6.7, 4.3, 0.6, color=ACCENT)
add_text(slide, 4.7, 6.75, 3.9, 0.5,
         "Total Feature Vector:  177 dimensions",
         size=18, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 8: LSTM Classifier
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 8, 0.6, "LSTM Classifier (Section 4.4)",
         size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.8, 0.95, 5)

add_card(slide, 0.8, 1.3, 6.0, 5.7)
tf = add_text(slide, 1.1, 1.5, 5.5, 0.5,
              "Model Architecture", size=22, color=ACCENT2, bold=True)
add_para(tf, "")
add_para(tf, "Input: 30 frames x 177 features", size=16, color=WHITE)
add_para(tf, "")
add_para(tf, "Bidirectional LSTM", size=18, color=GREEN, bold=True)
add_para(tf, "  • 2 layers, hidden size 128", size=15, color=WHITE)
add_para(tf, "  • Processes sequence forward AND backward", size=15, color=WHITE)
add_para(tf, "  • Captures context from past and future frames", size=15, color=WHITE)
add_para(tf, "")
add_para(tf, "Attention Mechanism", size=18, color=ORANGE, bold=True)
add_para(tf, "  • Learns to focus on the most important", size=15, color=WHITE)
add_para(tf, "    frames in the sequence", size=15, color=WHITE)
add_para(tf, "  • The actual fall frames get highest weight", size=15, color=WHITE)
add_para(tf, "")
add_para(tf, "Output: Binary classification", size=18, color=RED_SOFT, bold=True)
add_para(tf, "  • Probability: 0.0 (Normal) to 1.0 (Fall)", size=15, color=WHITE)
add_para(tf, "  • Threshold: 0.5", size=15, color=WHITE)
add_para(tf, "")
add_para(tf, "Total Parameters: 745,060", size=16, color=ACCENT)

add_card(slide, 7.2, 1.3, 5.3, 5.7)
tf = add_text(slide, 7.5, 1.5, 4.8, 0.5,
              "Why LSTM?", size=22, color=ACCENT2, bold=True)
add_para(tf, "")
add_para(tf, "A fall is a TEMPORAL event:", size=16, color=WHITE)
add_para(tf, "  Standing → Falling → On Ground", size=16, color=ORANGE)
add_para(tf, "")
add_para(tf, "You cannot detect a fall from a single", size=16, color=WHITE)
add_para(tf, "image — you need the SEQUENCE.", size=16, color=WHITE)
add_para(tf, "")
add_para(tf, "LSTM is designed for sequential data:", size=16, color=WHITE)
add_para(tf, "  • Remembers important past information", size=15, color=GREEN)
add_para(tf, "  • Forgets irrelevant information", size=15, color=GREEN)
add_para(tf, "  • Learns long-term dependencies", size=15, color=GREEN)
add_para(tf, "")
add_para(tf, "Training Details:", size=18, color=ORANGE, bold=True)
add_para(tf, "  • Optimizer: Adam (lr = 0.001)", size=15, color=WHITE)
add_para(tf, "  • Loss: Binary Cross-Entropy", size=15, color=WHITE)
add_para(tf, "  • Early stopping (patience = 10)", size=15, color=WHITE)
add_para(tf, "  • LR scheduling (reduce on plateau)", size=15, color=WHITE)
add_para(tf, "  • Gradient clipping (max norm = 1.0)", size=15, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 9: Project Structure
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 8, 0.6, "Project Structure",
         size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.8, 0.95, 4)

add_card(slide, 0.8, 1.3, 5.8, 5.7)
tf = add_text(slide, 1.1, 1.4, 5.3, 0.5,
              "Directory Layout", size=20, color=ACCENT2, bold=True)
add_para(tf, "")
add_para(tf, "fall-detection/", size=15, color=GREEN, bold=True)
add_para(tf, "  README.md", size=14, color=WHITE)
add_para(tf, "  config.py", size=14, color=WHITE)
add_para(tf, "  requirements.txt", size=14, color=WHITE)
add_para(tf, "  src/", size=14, color=ORANGE, bold=True)
add_para(tf, "    data/", size=14, color=ACCENT2)
add_para(tf, "      download_dataset.py", size=13, color=LIGHT_GRAY)
add_para(tf, "      preprocess.py", size=13, color=LIGHT_GRAY)
add_para(tf, "      dataset.py", size=13, color=LIGHT_GRAY)
add_para(tf, "    models/", size=14, color=ACCENT2)
add_para(tf, "      depth_estimator.py", size=13, color=LIGHT_GRAY)
add_para(tf, "      pose_estimator.py", size=13, color=LIGHT_GRAY)
add_para(tf, "      feature_extractor.py", size=13, color=LIGHT_GRAY)
add_para(tf, "      fall_detector.py", size=13, color=LIGHT_GRAY)
add_para(tf, "    train.py | evaluate.py | inference.py", size=13, color=LIGHT_GRAY)
add_para(tf, "  results/ | models/ | docs/", size=14, color=ACCENT2)

add_card(slide, 7.0, 1.3, 5.5, 5.7)
tf = add_text(slide, 7.3, 1.5, 5.0, 0.5,
              "Technologies Used", size=20, color=ACCENT2, bold=True)

techs = [
    ("Python 3.13", "Core language"),
    ("PyTorch", "Deep learning framework"),
    ("MiDaS", "Monocular depth estimation"),
    ("OpenCV", "Image/video processing"),
    ("NumPy / SciPy", "Numerical computation"),
    ("scikit-learn", "Evaluation metrics"),
    ("Matplotlib / Seaborn", "Visualization and plots"),
    ("Git / GitHub", "Version control"),
]
add_para(tf, "")
for tech, desc in techs:
    add_para(tf, f"  {tech}", size=16, color=GREEN, bold=True)
    add_para(tf, f"    {desc}", size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 10: Timeline
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 8, 0.6, "Six-Week Timeline",
         size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.8, 0.95, 4)

weeks = [
    ("Week 1", "Setup + Data",       "Environment, dataset download, Git repo", GREEN, True),
    ("Week 2", "Depth + Pose",       "MiDaS and pose estimation modules", GREEN, True),
    ("Week 3", "Features + Pipeline", "Feature engineering, preprocessing", GREEN, True),
    ("Week 4", "Training",           "LSTM model training, hyperparameter tuning", LIGHT_GRAY, False),
    ("Week 5", "Evaluation + Demo",  "Metrics, confusion matrix, demo video", LIGHT_GRAY, False),
    ("Week 6", "Report",             "Final documentation, presentation", LIGHT_GRAY, False),
]

for i, (week, focus, desc, clr, done) in enumerate(weeks):
    y = 1.2 + i * 0.95
    status_color = GREEN if done else BG_CARD
    border = GREEN if done else LIGHT_GRAY

    # Status indicator
    shape = add_card(slide, 0.8, y, 0.6, 0.7, color=status_color)
    if done:
        shape.text_frame.paragraphs[0].text = "DONE"
        shape.text_frame.paragraphs[0].font.size = Pt(10)
        shape.text_frame.paragraphs[0].font.color.rgb = BG_DARK
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_card(slide, 1.6, y, 10.9, 0.7, border=border)
    add_text(slide, 1.8, y + 0.05, 1.5, 0.3, week, size=16, color=clr, bold=True)
    add_text(slide, 3.5, y + 0.05, 2.5, 0.3, focus, size=16, color=WHITE, bold=True)
    add_text(slide, 6.3, y + 0.05, 6.0, 0.3, desc, size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 11: Next Steps
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 8, 0.6, "Next Steps (Weeks 4-6)",
         size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.8, 0.95, 5)

add_card(slide, 0.8, 1.3, 3.7, 5.5, border=ORANGE)
tf = add_text(slide, 1.1, 1.5, 3.2, 0.5,
              "Week 4: Training", size=20, color=ORANGE, bold=True)
add_para(tf, "")
add_para(tf, "• Train LSTM on full", size=15, color=WHITE)
add_para(tf, "  URFD dataset (70 seq)", size=15, color=WHITE)
add_para(tf, "• Hyperparameter tuning", size=15, color=WHITE)
add_para(tf, "  - Hidden size", size=14, color=LIGHT_GRAY)
add_para(tf, "  - Sequence length", size=14, color=LIGHT_GRAY)
add_para(tf, "  - Learning rate", size=14, color=LIGHT_GRAY)
add_para(tf, "• Cross-validation", size=15, color=WHITE)
add_para(tf, "• Log training curves", size=15, color=WHITE)

add_card(slide, 4.8, 1.3, 3.7, 5.5, border=RED_SOFT)
tf = add_text(slide, 5.1, 1.5, 3.2, 0.5,
              "Week 5: Evaluation", size=20, color=RED_SOFT, bold=True)
add_para(tf, "")
add_para(tf, "• Run evaluation metrics:", size=15, color=WHITE)
add_para(tf, "  - Accuracy", size=14, color=LIGHT_GRAY)
add_para(tf, "  - Sensitivity/Recall", size=14, color=LIGHT_GRAY)
add_para(tf, "  - Specificity", size=14, color=LIGHT_GRAY)
add_para(tf, "  - F1-Score, AUC-ROC", size=14, color=LIGHT_GRAY)
add_para(tf, "• Confusion matrix plot", size=15, color=WHITE)
add_para(tf, "• ROC curve plot", size=15, color=WHITE)
add_para(tf, "• Create demo video", size=15, color=WHITE)

add_card(slide, 8.8, 1.3, 3.7, 5.5, border=GREEN)
tf = add_text(slide, 9.1, 1.5, 3.2, 0.5,
              "Week 6: Report", size=20, color=GREEN, bold=True)
add_para(tf, "")
add_para(tf, "• Final results chapter", size=15, color=WHITE)
add_para(tf, "• Literature comparison", size=15, color=WHITE)
add_para(tf, "• Conclusion + future work", size=15, color=WHITE)
add_para(tf, "• Clean up repository", size=15, color=WHITE)
add_para(tf, "• Presentation slides", size=15, color=WHITE)
add_para(tf, "• Submit final report", size=15, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 12: Thank You
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_accent_line(slide, 3, 2.8, 7.3)
add_text(slide, 1, 3.0, 11.3, 1,
         "Thank You", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 4.2, 11.3, 0.5,
         "Questions & Discussion", size=24, color=ACCENT2, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.3, 11.3, 0.5,
         "GitHub: github.com/m25ai1127/fall-detection",
         size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── Save ──
output = os.path.join(os.path.dirname(__file__), "..", "docs", "presentation_weeks1_3.pptx")
os.makedirs(os.path.dirname(output), exist_ok=True)
prs.save(output)
print(f"Saved: {output}")
